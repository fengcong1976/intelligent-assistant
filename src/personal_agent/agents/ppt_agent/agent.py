# -*- coding: utf-8 -*-
"""
PPT 生成智能体
功能：根据用户提供的主题或大纲自动生成 PPT 文件
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os
import re
import logging
import json
import aiohttp
import asyncio
from typing import Optional, Dict, List, Any
from pathlib import Path

from personal_agent.agents.base import BaseAgent, Task

logger = logging.getLogger(__name__)


class PPTAgent(BaseAgent):
    """PPT 生成智能体"""
    
    LAYOUT_TITLE = 0
    LAYOUT_TITLE_CONTENT = 1
    LAYOUT_SECTION = 2
    LAYOUT_TWO_CONTENT = 3
    LAYOUT_COMPARISON = 4
    LAYOUT_BLANK = 6
    
    COLOR_THEMES = {
        'blue': {
            'primary': RGBColor(0, 112, 192),
            'secondary': RGBColor(68, 114, 196),
            'accent': RGBColor(91, 155, 213),
            'text': RGBColor(255, 255, 255),
            'dark_text': RGBColor(31, 73, 125),
        },
        'green': {
            'primary': RGBColor(0, 128, 0),
            'secondary': RGBColor(84, 130, 53),
            'accent': RGBColor(146, 208, 80),
            'text': RGBColor(255, 255, 255),
            'dark_text': RGBColor(56, 87, 35),
        },
        'red': {
            'primary': RGBColor(192, 0, 0),
            'secondary': RGBColor(196, 67, 68),
            'accent': RGBColor(255, 102, 102),
            'text': RGBColor(255, 255, 255),
            'dark_text': RGBColor(127, 0, 0),
        },
        'purple': {
            'primary': RGBColor(112, 48, 160),
            'secondary': RGBColor(128, 100, 162),
            'accent': RGBColor(180, 141, 208),
            'text': RGBColor(255, 255, 255),
            'dark_text': RGBColor(75, 0, 130),
        },
        'orange': {
            'primary': RGBColor(255, 153, 0),
            'secondary': RGBColor(237, 125, 49),
            'accent': RGBColor(255, 192, 0),
            'text': RGBColor(255, 255, 255),
            'dark_text': RGBColor(191, 96, 0),
        },
    }
    
    def __init__(self):
        super().__init__(name="ppt_agent", description="PPT 生成智能体 - 根据主题或大纲自动生成演示文稿")
        
        self.skill = {
            'name': 'ppt_agent',
            'description': 'PPT 生成智能体 - 根据主题或大纲自动生成演示文稿'
        }
        
        self._llm_gateway = None
        self._image_api_key = None
        
        self._register_capabilities()
    
    def _get_llm_gateway(self):
        """获取 LLM 网关"""
        if self._llm_gateway is None:
            from ...llm import LLMGateway
            from ...config import settings
            self._llm_gateway = LLMGateway(settings.llm)
        return self._llm_gateway
    
    def _get_image_api_key(self):
        """获取图片生成 API Key"""
        if self._image_api_key is None:
            from ...config import settings
            self._image_api_key = settings.llm.dashscope_api_key
        return self._image_api_key
    
    def _register_capabilities(self):
        """注册智能体能力"""
        self.register_capability(
            capability="create_ppt",
            description="创建PPT演示文稿。当用户要求创建PPT、制作演示文稿、生成幻灯片时调用此工具。",
            parameters={
                "type": "object",
                "properties": {
                    "title": {
                        "type": "string",
                        "description": "PPT标题"
                    },
                    "slides": {
                        "type": "array",
                        "description": "幻灯片内容列表",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string", "description": "幻灯片标题"},
                                "content": {"type": "string", "description": "幻灯片内容"}
                            }
                        }
                    },
                    "output_path": {
                        "type": "string",
                        "description": "输出文件路径（可选）"
                    }
                },
                "required": ["title"]
            },
            category="document",
            aliases=["ppt", "PPT", "演示文稿", "幻灯片"],
            alias_params={
                "创建ppt": {"title": "{arg1}"},
                "制作ppt": {"title": "{arg1}"},
                "生成ppt": {"title": "{arg1}"},
                "创建演示文稿": {"title": "{arg1}"},
                "制作演示文稿": {"title": "{arg1}"},
                "生成演示文稿": {"title": "{arg1}"}
            }
        )
        
        self.register_capability(
            capability="create_ppt_from_topic",
            description="根据主题自动生成PPT。当用户提供一个主题并要求生成PPT时调用此工具。",
            parameters={
                "type": "object",
                "properties": {
                    "topic": {
                        "type": "string",
                        "description": "PPT主题"
                    },
                    "num_slides": {
                        "type": "integer",
                        "description": "幻灯片数量（默认5页）",
                        "default": 5
                    },
                    "output_path": {
                        "type": "string",
                        "description": "输出文件路径（可选）"
                    }
                },
                "required": ["topic"]
            },
            category="document",
            aliases=["主题ppt", "关于ppt"],
            alias_params={
                "关于": {"topic": "{arg1}"},
                "主题": {"topic": "{arg1}"}
            }
        )
    
    async def execute_task(self, task: Task) -> str:
        """执行任务"""
        try:
            action = task.type
            
            if action == "create_ppt":
                title = task.params.get("title", "")
                slides = task.params.get("slides", [])
                output_path = task.params.get("output_path")
                
                if not slides:
                    return await self.create_ppt_from_topic(title, 5, output_path)
                
                return await self.create_ppt(title, slides, output_path)
            
            elif action == "create_ppt_from_topic":
                topic = task.params.get("topic", task.params.get("title", ""))
                num_slides = task.params.get("num_slides", 5)
                output_path = task.params.get("output_path")
                
                if not topic:
                    return "请提供 PPT 的主题"
                
                return await self.create_ppt_from_topic(topic, num_slides, output_path)
            
            else:
                return self.cannot_handle(f"未知操作: {action}")
                
        except Exception as e:
            error_msg = f"执行任务失败: {str(e)}"
            logger.error(error_msg)
            logger.exception("详细错误信息:")
            return error_msg
    
    def _get_help_info(self) -> str:
        """获取帮助信息"""
        return """PPT 生成智能体

功能说明
PPT 生成智能体可以根据您提供的主题、大纲或内容自动生成专业的演示文稿。

支持的操作
  - 创建 PPT：根据标题和内容创建 PPT 文件
  - 从主题生成：根据主题自动生成大纲并创建 PPT
  - 添加幻灯片：向现有 PPT 添加新页面

使用示例
  - "创建一个关于人工智能的 PPT" - 自动生成大纲并创建 PPT
  - "制作 PPT 标题：项目汇报 内容：概述、方案、计划" - 按指定内容创建
  - "生成 5 页关于气候变化的演示文稿" - 生成指定页数的 PPT
  - "创建一个产品介绍 PPT，包含产品特性、优势、价格" - 指定内容结构

输出格式
  - 文件格式：.pptx
  - 默认保存位置：文档目录（可在设置中配置）
  - 支持自定义保存路径"""

    async def create_ppt(
        self,
        title: str,
        slides: List[Dict[str, str]] = None,
        output_path: Optional[str] = None,
        subtitle: str = "AI 自动生成",
        theme: str = 'blue'
    ) -> str:
        """
        创建 PPT
        
        Args:
            title: PPT 标题
            slides: 幻灯片列表，每个元素包含 title 和 content
            output_path: 输出路径
            subtitle: 副标题
            theme: 颜色主题 (blue, green, red, purple, orange)
            
        Returns:
            生成结果
        """
        try:
            if slides is None:
                slides = []
            
            if output_path is None:
                safe_title = re.sub(r'[\\/:*?"<>|]', '_', title)
                from ...config import settings
                documents_dir = settings.directory.get_documents_dir()
                documents_dir.mkdir(parents=True, exist_ok=True)
                output_path = str(documents_dir / f"{safe_title}.pptx")
            
            colors = self.COLOR_THEMES.get(theme, self.COLOR_THEMES['blue'])
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            title_slide = prs.slides.add_slide(prs.slide_layouts[self.LAYOUT_BLANK])
            self._add_gradient_background(title_slide, colors)
            self._add_decorative_shapes(title_slide, colors)
            
            title_box = title_slide.shapes.add_textbox(
                Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.text = title
            title_para.font.size = Pt(54)
            title_para.font.bold = True
            title_para.font.color.rgb = colors['text']
            title_para.alignment = PP_ALIGN.CENTER
            
            subtitle_box = title_slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = colors['text']
            subtitle_para.alignment = PP_ALIGN.CENTER
            
            for i, slide_data in enumerate(slides):
                slide = prs.slides.add_slide(prs.slide_layouts[self.LAYOUT_BLANK])
                self._add_gradient_background(slide, colors, variant=i % 3)
                self._add_decorative_shapes(slide, colors)
                
                slide_title = slide_data.get('title', f'第 {i+1} 页')
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.9)
                )
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = slide_title
                title_para.font.size = Pt(36)
                title_para.font.bold = True
                title_para.font.color.rgb = colors['text']
                
                self._add_title_underline(slide, colors, Inches(0.5), Inches(1.25))
                
                content = slide_data.get('content', '')
                content_box = slide.shapes.add_textbox(
                    Inches(0.7), Inches(1.6), Inches(11.933), Inches(5.5)
                )
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                lines = [line.strip() for line in content.split('\n') if line.strip()]
                
                for j, line in enumerate(lines):
                    if j == 0:
                        para = content_frame.paragraphs[0]
                    else:
                        para = content_frame.add_paragraph()
                    
                    para.text = f"• {line}"
                    para.font.size = Pt(20)
                    para.font.color.rgb = RGBColor(50, 50, 50)
                    para.space_after = Pt(12)
                    para.level = 0
            
            prs.save(output_path)
            
            abs_path = os.path.abspath(output_path)
            logger.info(f"PPT 已生成：{abs_path}")
            
            return f"PPT 已生成：{abs_path}\n共 {len(slides) + 1} 页（含标题页）"
            
        except Exception as e:
            logger.error(f"生成 PPT 失败：{e}")
            import traceback
            traceback.print_exc()
            return f"生成 PPT 失败：{str(e)}"
    
    def _add_gradient_background(self, slide, colors: Dict, variant: int = 0):
        """添加渐变背景"""
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(7.5)
        )
        background.line.fill.background()
        
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 45 + variant * 15
        
        fill.gradient_stops[0].color.rgb = colors['primary']
        fill.gradient_stops[1].color.rgb = colors['secondary']
        
        spTree = slide.shapes._spTree
        sp = background._element
        spTree.remove(sp)
        spTree.insert(2, sp)
    
    def _add_decorative_shapes(self, slide, colors: Dict):
        """添加装饰形状"""
        shapes_data = [
            (MSO_SHAPE.OVAL, Inches(11), Inches(-0.5), Inches(3), Inches(3), 0.15),
            (MSO_SHAPE.OVAL, Inches(-0.5), Inches(5.5), Inches(2.5), Inches(2.5), 0.1),
            (MSO_SHAPE.RIGHT_TRIANGLE, Inches(10), Inches(5.5), Inches(4), Inches(2.5), 0.08),
        ]
        
        for shape_type, left, top, width, height, transparency in shapes_data:
            shape = slide.shapes.add_shape(shape_type, left, top, width, height)
            shape.line.fill.background()
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = colors['accent']
            
            spTree = slide.shapes._spTree
            sp = shape._element
            spTree.remove(sp)
            spTree.insert(3, sp)
    
    def _add_title_underline(self, slide, colors: Dict, left, top):
        """添加标题下划线"""
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top,
            Inches(2), Inches(0.05)
        )
        line.line.fill.background()
        line.fill.solid()
        line.fill.fore_color.rgb = colors['accent']
    
    async def create_ppt_from_topic(
        self,
        topic: str,
        num_slides: int = 5,
        output_path: Optional[str] = None
    ) -> str:
        """
        根据主题自动生成 PPT
        
        Args:
            topic: 主题
            num_slides: 幻灯片数量（不包括标题页）
            output_path: 输出路径
            
        Returns:
            生成结果
        """
        try:
            if output_path is None:
                safe_topic = re.sub(r'[\\/:*?"<>|]', '_', topic)
                from ...config import settings
                documents_dir = settings.directory.get_documents_dir()
                documents_dir.mkdir(parents=True, exist_ok=True)
                output_path = str(documents_dir / f"{safe_topic}.pptx")
            
            slides = await self._generate_outline_with_llm(topic, num_slides)
            
            return await self.create_ppt(topic, slides, output_path)
            
        except Exception as e:
            logger.error(f"从主题生成 PPT 失败：{e}")
            return f"生成 PPT 失败：{str(e)}"
    
    def _generate_outline(self, topic: str, num_slides: int) -> List[Dict[str, str]]:
        """
        生成 PPT 大纲（简单模板，可后续集成 AI）
        
        Args:
            topic: 主题
            num_slides: 幻灯片数量
            
        Returns:
            幻灯片列表
        """
        templates = {
            'default': [
                {'title': '概述', 'content': f'关于{topic}的基本介绍\n背景信息\n主要内容'},
                {'title': '核心内容', 'content': f'{topic}的核心要点\n关键概念\n重要特征'},
                {'title': '详细分析', 'content': f'{topic}的深入分析\n优势与特点\n应用场景'},
                {'title': '案例展示', 'content': f'{topic}的实际案例\n成功经验\n经验总结'},
                {'title': '总结与展望', 'content': f'{topic}的总结\n未来发展趋势\n行动建议'},
            ]
        }
        
        outline = templates['default']
        
        if num_slides <= len(outline):
            return outline[:num_slides]
        
        result = outline.copy()
        for i in range(len(outline), num_slides):
            result.append({
                'title': f'第{i + 1}部分',
                'content': f'关于{topic}的第{i + 1}部分内容\n要点一\n要点二\n要点三'
            })
        
        return result
    
    async def _generate_outline_with_llm(self, topic: str, num_slides: int) -> List[Dict[str, str]]:
        """
        使用 LLM 生成 PPT 大纲和内容
        
        Args:
            topic: 主题
            num_slides: 幻灯片数量
            
        Returns:
            幻灯片列表
        """
        llm = self._get_llm_gateway()
        if not llm:
            logger.warning("LLM 网关不可用，使用默认模板")
            return self._generate_outline(topic, num_slides)
        
        prompt = f"""你是一个专业的PPT内容策划师。请为以下主题生成一份详细、专业、内容丰富的PPT大纲和内容。

主题：{topic}
幻灯片数量：{num_slides} 页（不包括标题页）

请以JSON格式返回，格式如下：
{{
    "slides": [
        {{
            "title": "幻灯片标题",
            "content": "幻灯片内容（使用换行符分隔多个要点）"
        }}
    ]
}}

要求：
1. 每页幻灯片要有明确的主题和吸引人的标题
2. 内容要非常丰富、专业、有深度，避免泛泛而谈
3. 每页内容必须包含5-8个详细要点，用换行符分隔
4. 内容要有逻辑性，从概述到深入，最后总结
5. 确保内容与主题紧密相关，包含具体数据、案例或事实
6. 使用专业术语，但也要通俗易懂
7. 每个要点应该是一句完整的话，而不是简单的关键词

请直接返回JSON，不要包含其他文字。"""

        try:
            from ...llm import Message
            messages = [Message(role="user", content=prompt)]
            response = await llm.chat(messages)
            
            if response and response.content:
                content = response.content.strip()
                
                if content.startswith("```"):
                    lines = content.split("\n")
                    content = "\n".join(lines[1:-1] if lines[-1] == "```" else lines[1:])
                
                data = json.loads(content)
                slides = data.get("slides", [])
                
                if slides:
                    logger.info(f"LLM 生成了 {len(slides)} 页幻灯片内容")
                    return slides[:num_slides]
                    
        except json.JSONDecodeError as e:
            logger.error(f"解析 LLM 返回的 JSON 失败: {e}")
        except Exception as e:
            logger.error(f"LLM 生成大纲失败: {e}")
        
        return self._generate_outline(topic, num_slides)
    
    async def handle_intent(self, intent: Dict[str, Any]) -> str:
        """
        处理意图
        
        Args:
            intent: 意图字典
            
        Returns:
            处理结果
        """
        action = intent.get('action', '')
        params = intent.get('params', {})
        
        if action == 'create_ppt':
            title = params.get('title', '未命名演示文稿')
            slides = params.get('slides', [])
            output_path = params.get('output_path')
            
            if not slides:
                return await self.create_ppt_from_topic(title, 5, output_path)
            
            return await self.create_ppt(title, slides, output_path)
        
        elif action == 'create_ppt_from_topic':
            topic = params.get('topic', '')
            num_slides = params.get('num_slides', 5)
            output_path = params.get('output_path')
            
            if not topic:
                return "请提供 PPT 的主题"
            
            return await self.create_ppt_from_topic(topic, num_slides, output_path)
        
        return f"未知操作：{action}"
    
    async def process(self, text: str) -> str:
        """
        处理用户输入
        
        Args:
            text: 用户输入文本
            
        Returns:
            处理结果
        """
        text_lower = text.lower()
        
        if 'ppt' in text_lower or '演示文稿' in text:
            topic_match = re.search(r'关于(.+?)的', text)
            if topic_match:
                topic = topic_match.group(1).strip()
                num_match = re.search(r'(\d+)\s*页', text)
                num_slides = int(num_match.group(1)) if num_match else 5
                return await self.create_ppt_from_topic(topic, num_slides)
            
            title_match = re.search(r'标题[：:]\s*(.+?)(?:\s+内容|$)', text)
            content_match = re.search(r'内容[：:]\s*(.+)$', text)
            
            if title_match:
                title = title_match.group(1).strip()
                
                if content_match:
                    content_str = content_match.group(1).strip()
                    items = [item.strip() for item in re.split(r'[,，、\s]+', content_str) if item.strip()]
                    slides = [{'title': item, 'content': f'{item}的相关内容'} for item in items]
                else:
                    slides = []
                
                return await self.create_ppt(title, slides)
            
            create_patterns = [
                r'创建(?:一个)?(.+?)的?\s*ppt',
                r'制作(?:一个)?(.+?)的?\s*ppt',
                r'生成(?:一个)?(.+?)的?\s*ppt',
            ]
            
            for pattern in create_patterns:
                match = re.search(pattern, text_lower)
                if match:
                    topic = match.group(1).strip()
                    if topic:
                        return await self.create_ppt_from_topic(topic, 5)
        
        return "我可以帮您创建 PPT，请告诉我主题，例如：'创建一个关于人工智能的 PPT'"


def get_agent():
    """获取 PPT 智能体实例"""
    return PPTAgent()
