# -*- coding: utf-8 -*-
"""
PPT 生成智能体
功能：根据用户提供的主题或大纲自动生成 PPT 文件
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import re
import logging
from typing import Optional, Dict, List, Any

logger = logging.getLogger(__name__)


class PPTAgent:
    """PPT 生成智能体"""
    
    LAYOUT_TITLE = 0
    LAYOUT_TITLE_CONTENT = 1
    LAYOUT_SECTION = 2
    LAYOUT_TWO_CONTENT = 3
    LAYOUT_COMPARISON = 4
    LAYOUT_BLANK = 6
    
    def __init__(self):
        self.name = "ppt_agent"
        self.skill = {
            'name': 'ppt_agent',
            'description': 'PPT 生成智能体 - 根据主题或大纲自动生成演示文稿'
        }
        
        self.alias_params = {
            'create_ppt': ['创建ppt', '制作ppt', '生成ppt', '创建演示文稿', '制作演示文稿', '生成演示文稿'],
            'create_ppt_from_topic': ['关于', '主题'],
        }
    
    def _get_help_info(self) -> str:
        """获取帮助信息"""
        return """PPT 生成智能体

功能说明
PPT 生成智能体可以根据您提供的主题、大纲或内容自动生成专业的演示文稿。

支持的操作
  - 创建 PPT：根据标题和内容创建 PPT 文件
  - 从主题生成：根据主题自动生成大纲并创建 PPT
  - 添加幻灯片：向现有 PPT 添加新页面

使用示例
  - "创建一个关于人工智能的 PPT" - 自动生成大纲并创建 PPT
  - "制作 PPT 标题：项目汇报 内容：概述、方案、计划" - 按指定内容创建
  - "生成 5 页关于气候变化的演示文稿" - 生成指定页数的 PPT
  - "创建一个产品介绍 PPT，包含产品特性、优势、价格" - 指定内容结构

输出格式
  - 文件格式：.pptx
  - 默认保存位置：当前目录
  - 支持自定义保存路径"""

    async def create_ppt(
        self,
        title: str,
        slides: List[Dict[str, str]],
        output_path: Optional[str] = None,
        subtitle: str = "AI 自动生成"
    ) -> str:
        """
        创建 PPT
        
        Args:
            title: PPT 标题
            slides: 幻灯片列表，每个元素包含 title 和 content
            output_path: 输出路径
            subtitle: 副标题
            
        Returns:
            生成结果
        """
        try:
            if output_path is None:
                safe_title = re.sub(r'[\\/:*?"<>|]', '_', title)
                output_path = f"{safe_title}.pptx"
            
            prs = Presentation()
            
            slide = prs.slides.add_slide(prs.slide_layouts[self.LAYOUT_TITLE])
            slide.shapes.title.text = title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle
            
            for slide_data in slides:
                slide = prs.slides.add_slide(prs.slide_layouts[self.LAYOUT_TITLE_CONTENT])
                
                if 'title' in slide_data:
                    slide.shapes.title.text = slide_data['title']
                
                if 'content' in slide_data and len(slide.placeholders) > 1:
                    tf = slide.placeholders[1].text_frame
                    tf.text = slide_data['content']
            
            prs.save(output_path)
            
            abs_path = os.path.abspath(output_path)
            logger.info(f"PPT 已生成：{abs_path}")
            
            return f"PPT 已生成：{abs_path}\n共 {len(slides) + 1} 页（含标题页）"
            
        except Exception as e:
            logger.error(f"生成 PPT 失败：{e}")
            return f"生成 PPT 失败：{str(e)}"
    
    async def create_ppt_from_topic(
        self,
        topic: str,
        num_slides: int = 5,
        output_path: Optional[str] = None
    ) -> str:
        """
        根据主题自动生成 PPT
        
        Args:
            topic: 主题
            num_slides: 幻灯片数量（不包括标题页）
            output_path: 输出路径
            
        Returns:
            生成结果
        """
        try:
            if output_path is None:
                safe_topic = re.sub(r'[\\/:*?"<>|]', '_', topic)
                output_path = f"{safe_topic}.pptx"
            
            slides = self._generate_outline(topic, num_slides)
            
            return await self.create_ppt(topic, slides, output_path)
            
        except Exception as e:
            logger.error(f"从主题生成 PPT 失败：{e}")
            return f"生成 PPT 失败：{str(e)}"
    
    def _generate_outline(self, topic: str, num_slides: int) -> List[Dict[str, str]]:
        """
        生成 PPT 大纲（简单模板，可后续集成 AI）
        
        Args:
            topic: 主题
            num_slides: 幻灯片数量
            
        Returns:
            幻灯片列表
        """
        templates = {
            'default': [
                {'title': '概述', 'content': f'关于{topic}的基本介绍\n背景信息\n主要内容'},
                {'title': '核心内容', 'content': f'{topic}的核心要点\n关键概念\n重要特征'},
                {'title': '详细分析', 'content': f'{topic}的深入分析\n优势与特点\n应用场景'},
                {'title': '案例展示', 'content': f'{topic}的实际案例\n成功经验\n经验总结'},
                {'title': '总结与展望', 'content': f'{topic}的总结\n未来发展趋势\n行动建议'},
            ]
        }
        
        outline = templates['default']
        
        if num_slides <= len(outline):
            return outline[:num_slides]
        
        result = outline.copy()
        for i in range(len(outline), num_slides):
            result.append({
                'title': f'第{i + 1}部分',
                'content': f'关于{topic}的第{i + 1}部分内容\n要点一\n要点二\n要点三'
            })
        
        return result
    
    async def handle_intent(self, intent: Dict[str, Any]) -> str:
        """
        处理意图
        
        Args:
            intent: 意图字典
            
        Returns:
            处理结果
        """
        action = intent.get('action', '')
        params = intent.get('params', {})
        
        if action == 'create_ppt':
            title = params.get('title', '未命名演示文稿')
            slides = params.get('slides', [])
            output_path = params.get('output_path')
            
            if not slides:
                return await self.create_ppt_from_topic(title, 5, output_path)
            
            return await self.create_ppt(title, slides, output_path)
        
        elif action == 'create_ppt_from_topic':
            topic = params.get('topic', '')
            num_slides = params.get('num_slides', 5)
            output_path = params.get('output_path')
            
            if not topic:
                return "请提供 PPT 的主题"
            
            return await self.create_ppt_from_topic(topic, num_slides, output_path)
        
        return f"未知操作：{action}"
    
    async def process(self, text: str) -> str:
        """
        处理用户输入
        
        Args:
            text: 用户输入文本
            
        Returns:
            处理结果
        """
        text_lower = text.lower()
        
        if 'ppt' in text_lower or '演示文稿' in text:
            topic_match = re.search(r'关于(.+?)的', text)
            if topic_match:
                topic = topic_match.group(1).strip()
                num_match = re.search(r'(\d+)\s*页', text)
                num_slides = int(num_match.group(1)) if num_match else 5
                return await self.create_ppt_from_topic(topic, num_slides)
            
            title_match = re.search(r'标题[：:]\s*(.+?)(?:\s+内容|$)', text)
            content_match = re.search(r'内容[：:]\s*(.+)$', text)
            
            if title_match:
                title = title_match.group(1).strip()
                
                if content_match:
                    content_str = content_match.group(1).strip()
                    items = [item.strip() for item in re.split(r'[,，、\s]+', content_str) if item.strip()]
                    slides = [{'title': item, 'content': f'{item}的相关内容'} for item in items]
                else:
                    slides = []
                
                return await self.create_ppt(title, slides)
            
            create_patterns = [
                r'创建(?:一个)?(.+?)的?\s*ppt',
                r'制作(?:一个)?(.+?)的?\s*ppt',
                r'生成(?:一个)?(.+?)的?\s*ppt',
            ]
            
            for pattern in create_patterns:
                match = re.search(pattern, text_lower)
                if match:
                    topic = match.group(1).strip()
                    if topic:
                        return await self.create_ppt_from_topic(topic, 5)
        
        return "我可以帮您创建 PPT，请告诉我主题，例如：'创建一个关于人工智能的 PPT'"


def get_agent():
    """获取 PPT 智能体实例"""
    return PPTAgent()
